import os
import io
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
from .image_ocr import run_ocr


def extract(file_path: str, file_type: str) -> str:
    """
    MS Office 문서(docx, xlsx, pptx) 텍스트 + 이미지 OCR 추출
    """
    text_blocks = []

    try:
        if file_type == ".docx":
            text_blocks.extend(_extract_docx(file_path))
        elif file_type == ".xlsx":
            text_blocks.extend(_extract_xlsx(file_path))
        elif file_type == ".pptx":
            text_blocks.extend(_extract_pptx(file_path))
        else:
            return f"[ERROR] Unsupported file type: {file_type}"
    except Exception as e:
        print(f"[ERROR] Failed to extract from {file_path}: {e}")
        text_blocks.append(f"[ERROR] Failed to extract content: {e}")

    if not text_blocks:
        text_blocks.append("[EMPTY] No extractable content found.")

    return "\n".join(text_blocks)


# ======================================
# DOCX (Word)
# ======================================
def _extract_docx(file_path):
    results = []
    try:
        doc = Document(file_path)
        # 텍스트 추출
        for para in doc.paragraphs:
            if para.text.strip():
                results.append(para.text)

        # 이미지 추출
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_data = rel.target_part.blob
                image = Image.open(io.BytesIO(image_data))
                ocr_text = run_ocr(image)
                results.append(f"\n[IMAGE_OCR]\n{ocr_text}")

    except Exception as e:
        results.append(f"[WARN] DOCX extract failed: {e}")
    return results


# ======================================
# XLSX (Excel)
# ======================================
def _extract_xlsx(file_path):
    results = []
    try:
        wb = load_workbook(file_path, data_only=True)
        for sheet in wb.worksheets:
            results.append(f"[SHEET] {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell])
                if row_text.strip():
                    results.append(row_text)
    except Exception as e:
        results.append(f"[WARN] XLSX text extract failed: {e}")

    # Excel은 이미지가 zip 구조 내에 있음
    try:
        media_dir = os.path.join(os.path.dirname(file_path), "xl", "media")
        if os.path.exists(media_dir):
            for img_name in os.listdir(media_dir):
                with open(os.path.join(media_dir, img_name), "rb") as img_file:
                    image = Image.open(io.BytesIO(img_file.read()))
                    ocr_text = run_ocr(image)
                    results.append(f"\n[IMAGE_OCR:{img_name}]\n{ocr_text}")
    except Exception as e:
        results.append(f"[WARN] XLSX image extract failed: {e}")

    return results


# ======================================
# PPTX (PowerPoint)
# ======================================
def _extract_pptx(file_path):
    results = []
    try:
        prs = Presentation(file_path)

        # 텍스트 추출
        for slide_idx, slide in enumerate(prs.slides, 1):
            results.append(f"[SLIDE {slide_idx}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    results.append(shape.text.strip())

        # 이미지 추출
        for rel in prs.part.rels.values():
            if "image" in rel.target_ref:
                image_data = rel.target_part.blob
                image = Image.open(io.BytesIO(image_data))
                ocr_text = run_ocr(image)
                results.append(f"\n[IMAGE_OCR]\n{ocr_text}")

    except Exception as e:
        results.append(f"[WARN] PPTX extract failed: {e}")

    return results
